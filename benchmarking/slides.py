import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

import plots as plts


def load_benchmark_config(config_path):
    try:
        with open(config_path, "r") as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"Error: Configuration file not found at {config_path}")
        return {}
    except json.JSONDecodeError:
        print(f"Error: Failed to decode JSON from {config_path}")
        return {}


def save_to_ppt(slides, ppt_filename, plots_folder):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only layout

        # Remove default title placeholder to avoid overlap
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:
                slide.shapes._spTree.remove(shape._element)

        # Add custom full-width title text box, left aligned
        title_box = slide.shapes.add_textbox(
            0, Inches(0.15), prs.slide_width, Inches(1)
        )
        title_tf = title_box.text_frame
        title_tf.clear()
        title_p = title_tf.paragraphs[0]
        title_p.text = slide_data["title"]
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.alignment = PP_ALIGN.LEFT

        if slide_data.get("subtitle"):
            left = Inches(0.5)
            top = Inches(1.1)
            width = prs.slide_width - Inches(1)
            height = Inches(0.6)
            subtitle_box = slide.shapes.add_textbox(left, top, width, height)
            tf = subtitle_box.text_frame
            tf.text = slide_data["subtitle"]
            tf.paragraphs[0].font.size = Pt(18)
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT

        img_path = os.path.join(plots_folder, slide_data["filename"])
        if os.path.isfile(img_path):
            left_img = Inches(0.5)
            top_img = Inches(1.9)
            max_width = prs.slide_width - Inches(1)
            max_height = prs.slide_height - Inches(2.1)
            slide.shapes.add_picture(
                img_path, left_img, top_img, width=max_width, height=max_height
            )
    ppt_out = os.path.join(plots_folder, ppt_filename)
    prs.save(ppt_out)
    print(f"\nPowerPoint file created: {ppt_out}")


def create_fastest_methods_table(df_hp, streams_order):
    rows = []
    for s in streams_order:
        sub = df_hp[df_hp["streams"] == s]
        if not sub.empty:
            fastest = sub.loc[sub["time_per_frame"].idxmin()]
            rows.append(
                [
                    s,
                    fastest["method"],
                    fastest["time_per_frame"],
                    fastest["fps"],
                    fastest["cpu"],
                ]
            )

    return pd.DataFrame(
        rows, columns=["Streams", "Method", "Time/Frame (ms)", "FPS", "CPU (%)"]
    )


def add_fastest_methods_slide(slides, df_hp, streams_order, plots_folder, config_list):
    if not config_list:
        return

    config = config_list[0]

    tbl_filename = config["filename"]
    highlighted_png = config["highlighted_filename"]

    tbl_fastest = create_fastest_methods_table(df_hp, streams_order)

    # used for confluence
    plts.pretty_table(
        tbl_fastest,
        tbl_filename,
        plots_folder,
    )

    plts.save_highlighted_table_as_png(
        tbl_fastest,
        os.path.join(plots_folder, highlighted_png),
    )

    slides.append(
        {
            "title": config["title"],
            "subtitle": config["subtitle"],
            "filename": highlighted_png,
        }
    )


def add_scaling_charts(slides, df_hp, plots_folder, config_list):
    for cfg in config_list:
        plts.plot_scaling(
            df_hp,
            cfg["metric"],
            cfg["title"],
            cfg["ylabel"],
            cfg["filename"],
            plots_folder,
        )
        slides.append(
            {
                "title": cfg["title"],
                "subtitle": cfg["subtitle"],
                "filename": cfg["filename"],
            }
        )


def add_grouped_bar_charts(slides, df_hp, plots_folder, config_list):
    for cfg in config_list:
        plts.plot_grouped_bar(
            df_hp,
            cfg["metric"],
            cfg["chart_title"],
            cfg["ylabel"],
            cfg["filename"],
            plots_folder,
        )
        slides.append(
            {
                "title": cfg["slide_title"],
                "subtitle": cfg["slide_subtitle"],
                "filename": cfg["filename"],
            }
        )


def add_section_header(slides, title, subtitle):
    """Add a section header slide."""
    slides.append(
        {
            "title": title,
            "subtitle": subtitle,
            "filename": "blank.png",  ## why its used?
        }
    )


def create_detailed_table(df_sub):
    tbl = df_sub[
        ["method", "time_per_frame", "fps", "cpu", "memory", "mvs", "frames"]
    ].copy()
    tbl.columns = [
        "Method",
        "Time/frame (ms)",
        "FPS",
        "CPU (%)",
        "Mem Δ KB",
        "Total MVs",
        "Frames",
    ]
    return tbl


def add_detailed_tables(slides, df_hp, streams_order, plots_folder, config_list):
    if not config_list:
        return

    config = config_list[0]

    for streams in streams_order:
        df_sub = df_hp[df_hp["streams"] == streams]
        tbl = create_detailed_table(df_sub)

        # Save tables
        tbl_filename = config["filename"].format(streams=streams)
        highlighted_png = config["highlighted_filename"].format(streams=streams)

        # needed for confluence
        plts.pretty_table(tbl, tbl_filename, plots_folder)

        plts.save_highlighted_table_as_png(
            tbl, os.path.join(plots_folder, highlighted_png)
        )

        slides.append(
            {
                "title": config["title"].format(streams=streams),
                "subtitle": config["subtitle"].format(streams=streams),
                "filename": highlighted_png,
            }
        )


def add_per_stream_metric_charts(
    slides, df_hp, streams_order, plots_folder, config_list
):
    """Add individual bar charts for each stream count and metric."""
    for streams in streams_order:
        df_sub = df_hp[df_hp["streams"] == streams]

        for cfg in config_list:
            filename = cfg["filename"].format(streams=streams)
            chart_title = cfg["chart_title"].format(streams=streams)
            slide_title = cfg["slide_title"].format(streams=streams)
            slide_subtitle = cfg["slide_subtitle"].format(streams=streams)

            plts.plot_metric(
                df_sub,
                cfg["metric"],
                chart_title,
                cfg["ylabel"],
                filename,
                plots_folder,
                cfg["colormap"],
            )

            slides.append(
                {
                    "title": slide_title,
                    "subtitle": slide_subtitle,
                    "filename": filename,
                }
            )


def produce_slides(df_hp, slides_config_path, file_name, plots_folder):
    config = load_benchmark_config(slides_config_path)
    if not config:
        print("Aborting slide generation due to missing or invalid config.")
        return

    slides = []
    streams_order = sorted(df_hp["streams"].unique())

    # 1. Fastest methods table
    add_fastest_methods_slide(
        slides, df_hp, streams_order, plots_folder, config.get("fastest_methods", [])
    )

    # 2. Scaling line charts
    add_scaling_charts(slides, df_hp, plots_folder, config.get("scaling_metrics", []))

    # 3. Grouped bar charts
    add_grouped_bar_charts(
        slides, df_hp, plots_folder, config.get("grouped_bar_metrics", [])
    )

    # 4. Section header for detailed tables
    add_section_header(slides, "Detailed Tables", "Full Per-Streams Benchmark Results")

    # 5. Detailed tables per stream count
    add_detailed_tables(
        slides, df_hp, streams_order, plots_folder, config.get("detailed_tables", [])
    )

    # 6. Individual bar charts per stream and metric
    add_per_stream_metric_charts(
        slides, df_hp, streams_order, plots_folder, config.get("per_stream_metrics", [])
    )

    save_to_ppt(slides, file_name, plots_folder)
